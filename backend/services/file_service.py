"""
file_service.py — File upload, text extraction, and Supabase Storage operations.
Text is extracted at upload time and stored in the database to avoid re-processing.
"""

import io
import uuid
import pdfplumber
import docx
from pptx import Presentation
from services.supabase_client import supabase


MIME_BY_EXT = {
    "pdf":  "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt":  "application/vnd.ms-powerpoint",
}


def get_mime_type(file_name: str) -> str:
    """Return the correct MIME type for a given filename based on extension."""
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    return MIME_BY_EXT.get(ext, "application/octet-stream")


def upload_to_storage(file_bytes: bytes, file_name: str, course_id: str) -> str:
    """
    Upload a file to the 'course-materials' bucket with a unique path.

    Returns the storage path. Raises on failure.
    """
    safe_name = file_name.replace("/", "_").replace("\\", "_")
    unique_prefix = uuid.uuid4().hex[:8]
    path = f"{course_id}/{unique_prefix}_{safe_name}"
    supabase.storage.from_("course-materials").upload(
        path,
        file_bytes,
        {"content-type": get_mime_type(safe_name), "upsert": "false"},
    )
    return path


def ensure_storage_bucket():
    """Create the course-materials bucket if it does not exist."""
    try:
        buckets = supabase.storage.list_buckets()
        names = [b.name if hasattr(b, "name") else b.get("name") for b in buckets]
        if "course-materials" not in names:
            supabase.storage.create_bucket(
                "course-materials",
                options={"public": False, "file_size_limit": 10 * 1024 * 1024},
            )
            print("[storage] course-materials bucket created.")
    except Exception as e:
        print(f"[storage] Bucket check failed: {e}")


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Extract all text from a PDF file using pdfplumber.

    Args:
        file_bytes: Raw PDF content as bytes.

    Returns:
        Concatenated text from all pages, or empty string on failure.
    """
    text_parts = []
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
    except Exception:
        pass
    return "\n".join(text_parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """
    Extract all text from a DOCX file using python-docx.

    Args:
        file_bytes: Raw DOCX content as bytes.

    Returns:
        Concatenated paragraph text, or empty string on failure.
    """
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract all text from a PPTX file using python-pptx.

    Args:
        file_bytes: Raw PPTX content as bytes.

    Returns:
        Concatenated slide text, or empty string on failure.
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
        return "\n".join(text_parts)
    except Exception:
        return ""


def get_all_course_text(course_id: str) -> str:
    """
    Retrieve and concatenate extracted_text for all materials in a course.

    Args:
        course_id: UUID of the course.

    Returns:
        A single string with all extracted course text, or empty string if none.
    """
    res = supabase.table("materials").select("extracted_text").eq(
        "course_id", course_id
    ).execute()
    texts = [row["extracted_text"] for row in res.data if row.get("extracted_text")]
    return "\n\n---\n\n".join(texts)
